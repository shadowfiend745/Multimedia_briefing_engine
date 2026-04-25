import textwrap
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
from moviepy import ImageClip, AudioFileClip


def rgb_color(color_name):
    """
    Convert a color name to an RGBColor object.
    Args: color_name (str): The name of the color to be converted.
    Returns: A tuple representing the RGB values of the specified color.
    """

    colors = {
        "white": (255, 255, 255),
        "black": (0, 0, 0),
        "red": (255, 0, 0),
        "green": (0, 255, 0),
        "blue": (0, 0, 255),
        "yellow": (255, 255, 0),
        "cyan": (0, 255, 255),
        "magenta": (255, 0, 255)
    }

    return colors.get(color_name.lower(), (0, 0, 0)) # default to black if color not found


def add_textbox(slide, left, top, width, height, text, font_size=24, bold=False, italic=False, color="white"):
    """
    Add a textbox to the given slide with the specified properties.
    Args: slide: The slide to which the textbox will be added.
          left: The left position of the textbox in inches.
          top: The top position of the textbox in inches.
          width: The width of the textbox in inches.
          height: The height of the textbox in inches.
          text (str): The text to be displayed in the textbox.
          font_size (int, optional): The font size of the text. Default is 24.
          bold (bool, optional): Whether the text should be bold. Default is False.
          italic (bool, optional): Whether the text should be italic. Default is False.
          color (str, optional): The color of the text. Default is "white".
    """

    r, g, b = rgb_color(color)
    # When we add a textbox to a slide, it creates a shape object that contains a text frame. 
    # We can access this text frame to set the properties of the text.
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    # when pptx creates a textbox, it automatically adds a paragraph to the text frame. 
    # We can access this paragraph and set its properties.
    p = text_frame.paragraphs[0] 
    p.text = text
    # The paragraph contains a run object that represents a contiguous run of text with the same formatting.
    run = p.runs[0] 
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor(r, g, b)


def parser(text):
    """
    Parse the briefing text file and extract the intro, key points, and summary.
    Args: text (str): The path to the briefing text file.
    Returns: intro (str): The introduction text.
             key_points (list of str): A list of key points.
             summary (str): The summary text.   
    """

    with open(text, "r", encoding='cp1252') as f:
        content = f.read()

    lines = content.strip().split('\n')
    intro = ''
    key_points = []
    summary = ''
    current_section = None

    # We iterate through each line of the text file and determine which section it belongs to (intro, key points, or summary) based on the section headers.
    for line in lines:
        if not line:
            continue
        if line == 'INTRO':
            current_section = 'intro'
        elif line == 'KEY POINTS':
            current_section = 'key_points'
        elif line == 'SUMMARY':
            current_section = 'summary'
        elif current_section == 'intro':
            intro = line
        elif current_section == 'key_points':
            key_point = line.lstrip('•–- ').strip()
            key_points.append(key_point)
        elif current_section == 'summary':
            summary = line
    
    return intro, key_points, summary


def slide_image(title, contents, filename, font_title, font_body, background_image, WIDTH, HEIGHT):
    """
    Create a slide image with the given title and contents.
    Args: title (str): The title of the slide.
          contents (list of str): The content lines to be added to the slide.
          filename (str): The name of the output image file.
          font_title (ImageFont): The font to be used for the title.
          font_body (ImageFont): The font to be used for the body text.
          background_image: (str or None): The path to the background image file. If None, a black background will be used.
          WIDTH (int): The width of the slide image.
          HEIGHT (int): The height of the slide image.
    """
    margin = 60
    title_Y = 50
    body_start_Y = 140
    line_spacing = 30
    paragraph_spacing = 10

    if background_image:
        # Use LANCZOS filter for high-quality downsampling when resizing the background image to fit the slide dimensions.
        img = Image.open(background_image).resize((WIDTH, HEIGHT), Image.LANCZOS) #
    else:
        img = Image.new('RGB', (WIDTH, HEIGHT), color=rgb_color("black"))

    # To calculate the maximum number of characters per line, 
    # we need to determine the average width of a character in the body font.
    bbox = font_body.getbbox("A")
    char_height = bbox[3] - bbox[1]
    average_char_width = font_body.getlength("abcdefghijklmnopqrstuvwxyz") / 26
    max_chars_per_line = int((WIDTH - 2 * margin) / average_char_width)

    # We use the ImageDraw module to draw the title and body text onto the slide image.
    draw = ImageDraw.Draw(img)
    draw.text((margin, title_Y), title, fill=rgb_color("white"), font=font_title)

    # We wrap the body text to fit within the slide width and draw each line of text onto the image, 
    # ensuring that we do not exceed the slide height.
    y_text = body_start_Y
    for content in contents:
        wrapped_text = textwrap.wrap(content, width=max_chars_per_line)
        # print(type(wrapped_text))
        for line in wrapped_text:
            if y_text + char_height > HEIGHT - 40:
                print(f"file: {filename} - Content exceeds slide height")
                raise ValueError("Please reduce the amount of text or increase the slide height.")

            draw.text((margin, y_text), line, fill=rgb_color("white"), font=font_body)
            y_text += line_spacing
        y_text += paragraph_spacing

    img.save(filename)
    print(f"Slide saved as {filename}")


def clips_generator(image_files, audio_files):
    """
    Generate a video clip from the given image and audio files.
    Args: image_files (str): The path to the image file to be used as the video frame.
          audio_files (str): The path to the audio file to be used as the video audio.
          Returns: clip: A video clip object that can be concatenated with other clips to create the final video.
    """

    audio = AudioFileClip(audio_files)
    clip = ImageClip(image_files, duration=audio.duration).with_audio(audio)
    return clip